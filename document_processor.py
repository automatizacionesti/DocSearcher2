from pathlib import Path
from typing import List, Dict, Generator, Tuple
import pdfplumber
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
import io
import os
import re

class DocumentProcessor:
    def __init__(self):
        # Configura la ruta a Tesseract en Windows
        if os.name == 'nt':  # Windows
            pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
            # Agregar la ruta de poppler
            self.poppler_path = r'C:\Program Files\poppler\Library\bin'
        else:
            self.poppler_path = None
        
        self.processors = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.xlsx': self._process_xlsx,
            '.pptx': self._process_pptx,
            '.png': self._process_image,
            '.jpg': self._process_image,
            '.jpeg': self._process_image,
            '.tiff': self._process_image,
            '.bmp': self._process_image
        }
        self.progress_callback = None
        self.current_file_num = 0
        self.total_files = 0
    
    def count_processable_files(self, directory: Path) -> Tuple[int, List[Path]]:
        """Cuenta la cantidad de archivos procesables y retorna la lista de archivos."""
        processable_files = []
        for file_path in directory.rglob('*'):
            if file_path.suffix.lower() in self.processors:
                processable_files.append(file_path)
        return len(processable_files), processable_files

    def get_processable_files(self, directory: Path) -> Generator[Path, None, None]:
        """Genera una lista de archivos procesables en el directorio."""
        for file_path in directory.rglob('*'):
            if file_path.suffix.lower() in self.processors:
                yield file_path

    def get_file_type(self, file_path: Path) -> str:
        """Obtiene el tipo de archivo basado en su extensión."""
        ext = file_path.suffix.lower()
        if ext == '.pdf':
            return 'pdf'
        elif ext == '.docx':
            return 'docx'
        elif ext == '.xlsx':
            return 'xlsx'
        elif ext == '.pptx':
            return 'pptx'
        elif ext == '.txt':
            return 'txt'
        elif ext in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp']:
            return 'image'
        return 'others'

    def process_file(self, file_path: Path, keywords: List[str]) -> Dict:
        """Procesa un archivo individual y retorna los resultados."""
        try:
            total_steps = 4  # Total de pasos en el proceso
            current_step = 0
            
            # Notificar inicio de procesamiento del archivo
            if self.progress_callback:
                current_step += 1
                self.progress_callback(
                    current=self.current_file_num,
                    total=self.total_files,
                    file_path=file_path,
                    status=f"Iniciando procesamiento de {file_path.name}",
                    file_progress=(current_step / total_steps) * 100
                )
            
            processor = self.processors.get(file_path.suffix.lower())
            if not processor:
                return {
                    'file_name': str(file_path),
                    'file_type': self.get_file_type(file_path),
                    'error': 'Tipo de archivo no soportado'
                }

            # Notificar progreso de extracción de texto
            if self.progress_callback:
                current_step += 1
                self.progress_callback(
                    current=self.current_file_num,
                    total=self.total_files,
                    file_path=file_path,
                    status=f"Extrayendo texto de {file_path.name}",
                    file_progress=(current_step / total_steps) * 100
                )

            text = processor(file_path)
            if isinstance(text, dict) and 'error' in text:
                return text
            
            if not text:
                return {
                    'file_name': str(file_path),
                    'file_type': self.get_file_type(file_path),
                    'error': 'No se pudo extraer texto del archivo'
                }

            # Notificar inicio de búsqueda de palabras clave
            if self.progress_callback:
                current_step += 1
                self.progress_callback(
                    current=self.current_file_num,
                    total=self.total_files,
                    file_path=file_path,
                    status=f"Buscando coincidencias en {file_path.name}",
                    file_progress=(current_step / total_steps) * 100
                )

            findings = []
            text_lower = text.lower()
            
            for keyword in keywords:
                keyword_lower = keyword.lower()
                start = 0
                while True:
                    pos = text_lower.find(keyword_lower, start)
                    if pos == -1:
                        break
                    
                    context = self._get_context(text, keyword)
                    if context:
                        findings.append({
                            'path': str(file_path),
                            'file_name': file_path.name,
                            'palabra': keyword,
                            'contexto': context,
                            'file_type': self.get_file_type(file_path)
                        })
                    start = pos + 1

            # Notificar finalización del procesamiento
            if self.progress_callback:
                current_step += 1
                self.progress_callback(
                    current=self.current_file_num,
                    total=self.total_files,
                    file_path=file_path,
                    status=f"Archivo {file_path.name} procesado completamente",
                    file_progress=(current_step / total_steps) * 100
                )

            if findings:
                return {
                    'file_name': file_path.name,
                    'file_type': self.get_file_type(file_path),
                    'findings': findings
                }
            return {
                'file_name': file_path.name,
                'file_type': self.get_file_type(file_path),
                'findings': []
            }

        except Exception as e:
            return {
                'file_name': str(file_path),
                'file_type': self.get_file_type(file_path),
                'error': f"Error procesando archivo: {str(e)}"
            }

    def process_directory(self, directory: Path, keywords: List[str], progress_callback=None) -> Generator:
        """Procesa todos los archivos en un directorio buscando las palabras clave."""
        try:
            print(f"Iniciando búsqueda en directorio: {directory}")
            print(f"Palabras clave: {keywords}")
            
            # Obtener lista de archivos procesables
            files = list(self.get_processable_files(directory))
            total_files = len(files)
            
            print(f"Total de archivos a procesar: {total_files}")
            
            for current, file_path in enumerate(files, 1):
                try:
                    print(f"Procesando archivo {current}/{total_files}: {file_path}")
                    
                    # Actualizar progreso
                    if progress_callback:
                        progress_callback(current, total_files, file_path, "Procesando archivo...")

                    # Procesar archivo
                    result = self.process_file(file_path, keywords)
                    
                    if result:
                        print(f"Encontradas coincidencias en {file_path}: {len(result.get('findings', []))}")
                        yield result

                except Exception as e:
                    print(f"Error procesando archivo {file_path}: {str(e)}")
                    yield {
                        'error': str(e),
                        'file_name': str(file_path)
                    }
            
        except Exception as e:
            print(f"Error general en process_directory: {str(e)}")
            if progress_callback:
                progress_callback(0, 0, None, f"Error: {str(e)}")
    
    def _process_file(self, file_path: Path, keywords: List[str]) -> Generator[Dict, None, None]:
        """Procesa un archivo individual y busca las palabras o frases clave."""
        processor = self.processors.get(file_path.suffix.lower())
        if not processor:
            return

        try:
            text = processor(file_path)
            text_lower = text.lower()
            
            for keyword in keywords:
                keyword_lower = keyword.lower()
                
                # Buscar todas las ocurrencias de la palabra/frase
                start = 0
                while True:
                    pos = text_lower.find(keyword_lower, start)
                    if pos == -1:
                        break
                        
                    context = self._get_context(text, keyword)
                    if context:
                        yield {
                            'path': str(file_path),
                            'file_name': file_path.name,
                            'palabra': keyword,
                            'contexto': context,
                            'file_type': file_path.suffix.lstrip('.')
                        }
                    
                    start = pos + 1

        except Exception as e:
            print(f"Error procesando {file_path}: {str(e)}")
    
    def _get_context(self, text: str, search_term: str, context_size: int = 100) -> str:
        """
        Obtiene el contexto alrededor de una palabra o frase encontrada.
        Ahora soporta búsqueda de frases completas.
        """
        try:
            # Buscar la posición de la palabra/frase (ignorando mayúsculas/minúsculas)
            text_lower = text.lower()
            term_lower = search_term.lower()
            start_pos = text_lower.find(term_lower)
            
            if start_pos == -1:
                return ""
            
            # Calcular el inicio y fin del contexto
            context_start = max(0, start_pos - context_size)
            context_end = min(len(text), start_pos + len(search_term) + context_size)
            
            # Extraer el contexto
            context = text[context_start:context_end]
            
            # Resaltar la palabra/frase encontrada
            original_term = text[start_pos:start_pos + len(search_term)]
            context = context.replace(original_term, f"**{original_term}**")
            
            # Agregar elipsis si el contexto está truncado
            if context_start > 0:
                context = "..." + context
            if context_end < len(text):
                context = context + "..."
                
            return context
        except Exception as e:
            print(f"Error obteniendo contexto: {str(e)}")
            return ""
    
    def _process_pdf(self, file_path: Path) -> str:
        """Extrae texto de archivos PDF usando tanto extracción directa como OCR."""
        try:
            with pdfplumber.open(file_path) as pdf:
                text = ""
                total_pages = len(pdf.pages)
                
                for i, page in enumerate(pdf.pages, 1):
                    try:
                        page_text = page.extract_text()
                        if not page_text:
                            # Notificar que se inicia proceso OCR
                            if self.progress_callback:
                                self.progress_callback(
                                    current=self.current_file_num,
                                    total=self.total_files,
                                    file_path=file_path,
                                    status=f"Iniciando OCR en página {i}/{total_pages}",
                                    file_progress=0
                                )
                            
                            images = convert_from_path(file_path, poppler_path=self.poppler_path)
                            
                            for image in images:
                                # Actualizar progreso durante OCR
                                if self.progress_callback:
                                    self.progress_callback(
                                        current=self.current_file_num,
                                        total=self.total_files,
                                        file_path=file_path,
                                        status=f"Procesando OCR en página {i}/{total_pages}",
                                        file_progress=50
                                    )
                                
                                text += pytesseract.image_to_string(image, lang='spa')
                        
                        text += page_text + "\n"
                        
                        # Actualizar progreso después de cada página
                        if self.progress_callback:
                            progress = (i / total_pages) * 100
                            self.progress_callback(
                                current=self.current_file_num,
                                total=self.total_files,
                                file_path=file_path,
                                status=f"Página {i}/{total_pages} completada",
                                file_progress=progress
                            )
                        
                    except Exception as e:
                        print(f"Error en página {i}: {str(e)}")
                        continue
                
                return text
                
        except Exception as e:
            error_msg = f"Error procesando PDF {file_path.name}: {str(e)}"
            if self.progress_callback:
                self.progress_callback(
                    current=self.current_file_num,
                    total=self.total_files,
                    file_path=file_path,
                    status=error_msg,
                    file_progress=0
                )
            return {
                'file_name': str(file_path),
                'file_type': 'pdf',
                'error': error_msg
            }

    def _process_image(self, file_path: Path) -> str:
        """Procesa imágenes usando OCR."""
        try:
            with open(file_path, 'rb') as image_file:
                return self._apply_ocr(image_file.read())
        except Exception as e:
            print(f"Error procesando imagen {file_path}: {str(e)}")
            return ""

    def _apply_ocr(self, image_data: bytes) -> str:
        """Aplica OCR a una imagen."""
        try:
            # Convertir bytes a imagen
            image = Image.open(io.BytesIO(image_data))
            
            # Aplicar OCR
            text = pytesseract.image_to_string(image, lang='spa+eng')
            return text.strip()
        except Exception as e:
            print(f"Error en OCR: {str(e)}")
            return ""
    
    def _process_docx(self, file_path: Path) -> str:
        """Extrae texto de archivos Word."""
        doc = Document(file_path)
        return '\n'.join(paragraph.text for paragraph in doc.paragraphs)
    
    def _process_xlsx(self, file_path: Path) -> str:
        """Extrae texto de archivos Excel."""
        wb = load_workbook(file_path, read_only=True)
        text = []
        for sheet in wb.worksheets:
            for row in sheet.rows:
                text.extend(str(cell.value or '') for cell in row)
        return '\n'.join(text)
    
    def _process_pptx(self, file_path: Path) -> str:
        """Extrae texto de archivos PowerPoint."""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)

    def process_image_file(self, file_path: Path, keywords: List[str], progress_callback=None) -> Dict:
        """Procesa archivos de imagen usando OCR."""
        try:
            # Notificar inicio de proceso OCR
            if progress_callback:
                progress_callback(0, 1, file_path, "Iniciando proceso OCR en imagen", 0)
            
            # Usar pytesseract para extraer texto
            image = Image.open(file_path)
            
            # Notificar progreso OCR
            if progress_callback:
                progress_callback(0, 1, file_path, "Procesando OCR: Convirtiendo imagen a texto", 50)
            
            text = pytesseract.image_to_string(image, lang='spa')
            
            # Notificar finalización OCR
            if progress_callback:
                progress_callback(0, 1, file_path, "Finalizando proceso OCR", 100)
            
            return self._process_text(text, keywords, file_path)
            
        except Exception as e:
            return {
                'error': f"Error procesando imagen: {str(e)}",
                'file_name': file_path.name
            }

    def process_pdf_file(self, file_path: Path, keywords: List[str], progress_callback=None) -> Dict:
        """Procesa archivos PDF."""
        try:
            with pdfplumber.open(file_path) as pdf:
                text = ""
                total_pages = len(pdf.pages)
                
                for i, page in enumerate(pdf.pages, 1):
                    try:
                        page_text = page.extract_text()
                        if not page_text:
                            # Notificar inicio de OCR en página PDF
                            if progress_callback:
                                progress_callback(
                                    i, total_pages, 
                                    file_path,
                                    f"Iniciando OCR en página {i}/{total_pages}",
                                    0
                                )
                            
                            # Si no hay texto extraíble, intentar OCR
                            images = convert_from_path(
                                file_path,
                                first_page=i,
                                last_page=i
                            )
                            
                            # Notificar progreso OCR
                            if progress_callback:
                                progress_callback(
                                    i, total_pages,
                                    file_path,
                                    f"Procesando OCR en página {i}/{total_pages}",
                                    50
                                )
                            
                            for image in images:
                                page_text += pytesseract.image_to_string(image, lang='spa')
                        
                        text += page_text + "\n"
                        
                        if progress_callback:
                            progress = (i / total_pages) * 100
                            progress_callback(i, total_pages, file_path, f"Procesando página {i}/{total_pages}", progress)
                            
                    except Exception as e:
                        print(f"Error en página {i}: {str(e)}")
                        continue
                
                return self._process_text(text, keywords, file_path)
                
        except Exception as e:
            return {
                'error': f"Error procesando PDF: {str(e)}",
                'file_name': file_path.name
            } 